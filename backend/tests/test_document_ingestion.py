import os
import sys
import tempfile
import unittest

from pypdf import PdfWriter

from app.document_ingestion import (
    ChunkPolicyEngine,
    DoclingParserAdapter,
    DocumentParserRegistry,
    MinerUPdfParserAdapter,
)
from app.document_models import CanonicalBlock, KnowledgeChunkRecord


class DocumentIngestionModelTestCase(unittest.TestCase):
    @staticmethod
    def _build_valid_pdf_bytes() -> bytes:
        writer = PdfWriter()
        writer.add_blank_page(width=300, height=300)
        output = tempfile.SpooledTemporaryFile()
        writer.write(output)
        output.seek(0)
        payload = output.read()
        output.close()
        return payload

    def test_chunk_record_keeps_position_metadata(self) -> None:
        chunk = KnowledgeChunkRecord(
            id="chunk-1",
            source_file_id="file-1",
            organization_id="org-1",
            workspace_id="workspace-finance",
            scope="workspace",
            title="报销制度",
            content="报销审批正文",
            block_type="paragraph",
            section_path=["第一章", "审批流程"],
            page_number=2,
            sheet_name=None,
            slide_number=None,
            chunk_index=0,
            token_count_estimate=32,
            metadata={"source_type": "pdf"},
        )

        self.assertEqual(chunk.page_number, 2)
        self.assertEqual(chunk.section_path, ["第一章", "审批流程"])

    def test_canonical_block_can_carry_sheet_and_slide_metadata(self) -> None:
        block = CanonicalBlock(
            block_type="sheet",
            text="一月报销数据",
            section_path=["财务报表"],
            page_number=None,
            sheet_name="January",
            slide_number=None,
            metadata={"source_type": "xlsx"},
        )

        self.assertEqual(block.sheet_name, "January")
        self.assertEqual(block.metadata["source_type"], "xlsx")

    def test_docling_adapter_returns_canonical_blocks(self) -> None:
        adapter = DoclingParserAdapter()
        blocks = adapter._normalize_mock_blocks(
            [
                {"type": "heading", "text": "第一章 总则", "page_number": 1},
                {"type": "paragraph", "text": "这里是正文", "page_number": 1},
            ]
        )

        self.assertEqual(blocks[0].block_type, "heading")
        self.assertEqual(blocks[1].text, "这里是正文")

    def test_docling_adapter_can_parse_plain_text_as_single_paragraph(self) -> None:
        adapter = DoclingParserAdapter()
        with tempfile.NamedTemporaryFile("w", suffix=".txt", delete=False, encoding="utf-8") as handle:
            handle.write("报销制度正文")
            path = handle.name

        try:
            blocks = adapter.parse_file(path, "txt")
        finally:
            os.remove(path)

        self.assertEqual(len(blocks), 1)
        self.assertEqual(blocks[0].block_type, "paragraph")
        self.assertEqual(blocks[0].text, "报销制度正文")

    def test_docling_adapter_strips_utf8_bom_from_plain_text(self) -> None:
        adapter = DoclingParserAdapter()
        with tempfile.NamedTemporaryFile("w", suffix=".txt", delete=False, encoding="utf-8-sig") as handle:
            handle.write("差旅审批需要二次确认")
            path = handle.name

        try:
            blocks = adapter.parse_file(path, "txt")
        finally:
            os.remove(path)

        self.assertEqual(blocks[0].text, "差旅审批需要二次确认")
        self.assertFalse(blocks[0].text.startswith("\ufeff"))

    def test_docling_adapter_surfaces_pdf_runtime_error_when_mineru_is_missing(self) -> None:
        adapter = DoclingParserAdapter()

        with self.assertRaises(ValueError) as context:
            adapter.parse_file("fake.pdf", "pdf")

        self.assertEqual(str(context.exception), "mineru_executable_not_found")

    def test_parser_registry_can_parse_docx_headings_and_paragraphs(self) -> None:
        from docx import Document

        registry = DocumentParserRegistry()
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as handle:
            path = handle.name

        try:
            document = Document()
            document.add_heading("财务制度", level=1)
            document.add_paragraph("报销审批需要部门负责人确认。")
            document.save(path)

            blocks = registry.parse_file(path, "docx")
        finally:
            os.remove(path)

        self.assertTrue(any(block.block_type == "heading" for block in blocks))
        self.assertTrue(any(block.block_type == "paragraph" for block in blocks))
        self.assertEqual(blocks[0].text, "财务制度")

    def test_parser_registry_keeps_slide_boundary_for_pptx(self) -> None:
        from pptx import Presentation

        registry = DocumentParserRegistry()
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as handle:
            path = handle.name

        try:
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "销售策略"
            slide.placeholders[1].text = "价格策略需要审批。"
            presentation.save(path)

            blocks = registry.parse_file(path, "pptx")
        finally:
            os.remove(path)

        self.assertTrue(all(block.slide_number == 1 for block in blocks))
        self.assertTrue(any(block.block_type == "heading" for block in blocks))
        self.assertTrue(any(block.block_type in {"slide", "paragraph"} for block in blocks))

    def test_parser_registry_keeps_sheet_name_for_xlsx(self) -> None:
        from openpyxl import Workbook

        registry = DocumentParserRegistry()
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as handle:
            path = handle.name

        try:
            workbook = Workbook()
            worksheet = workbook.active
            worksheet.title = "January"
            worksheet.append(["项目", "金额"])
            worksheet.append(["差旅", 200])
            workbook.save(path)

            blocks = registry.parse_file(path, "xlsx")
        finally:
            os.remove(path)

        self.assertTrue(all(block.sheet_name == "January" for block in blocks))
        self.assertTrue(any(block.block_type in {"table", "sheet"} for block in blocks))

    def test_mineru_pdf_adapter_maps_cli_markdown_to_canonical_blocks(self) -> None:
        adapter_script = """
from pathlib import Path
import sys

args = sys.argv[1:]
pdf_path = Path(args[args.index("-p") + 1])
output_dir = Path(args[args.index("-o") + 1])
nested_dir = output_dir / pdf_path.stem / "auto"
nested_dir.mkdir(parents=True, exist_ok=True)
(nested_dir / f"{pdf_path.stem}.md").write_text(
    "# 财务审批\\n\\n报销审批需要部门负责人确认。\\n\\n| 项目 | 金额 |\\n| --- | --- |\\n| 差旅 | 200 |\\n",
    encoding="utf-8",
)
"""
        with tempfile.NamedTemporaryFile("w", suffix=".py", delete=False, encoding="utf-8") as script_handle:
            script_handle.write(adapter_script)
            script_path = script_handle.name

        with tempfile.NamedTemporaryFile("wb", suffix=".pdf", delete=False) as pdf_handle:
            pdf_handle.write(b"%PDF-1.4 fake")
            pdf_path = pdf_handle.name

        try:
            adapter = MinerUPdfParserAdapter(
                executable=sys.executable,
                extra_args=[script_path],
            )
            blocks = adapter.parse_file(pdf_path)
        finally:
            os.remove(script_path)
            os.remove(pdf_path)

        self.assertEqual(blocks[0].block_type, "heading")
        self.assertEqual(blocks[0].text, "财务审批")
        self.assertTrue(any(block.block_type == "paragraph" for block in blocks))
        self.assertTrue(any(block.block_type == "table" for block in blocks))

    def test_mineru_pdf_adapter_passes_backend_and_lang_arguments(self) -> None:
        adapter_script = """
from pathlib import Path
import sys

args = sys.argv[1:]
pdf_path = Path(args[args.index("-p") + 1])
output_dir = Path(args[args.index("-o") + 1])
backend = args[args.index("-b") + 1]
lang = args[args.index("-l") + 1]
method = args[args.index("-m") + 1]
nested_dir = output_dir / pdf_path.stem / "auto"
nested_dir.mkdir(parents=True, exist_ok=True)
(nested_dir / f"{pdf_path.stem}.md").write_text(
    f"# backend={backend}\\n\\nlang={lang}\\n\\nmethod={method}",
    encoding="utf-8",
)
"""
        with tempfile.NamedTemporaryFile("w", suffix=".py", delete=False, encoding="utf-8") as script_handle:
            script_handle.write(adapter_script)
            script_path = script_handle.name

        with tempfile.NamedTemporaryFile("wb", suffix=".pdf", delete=False) as pdf_handle:
            pdf_handle.write(b"%PDF-1.4 fake")
            pdf_path = pdf_handle.name

        try:
            adapter = MinerUPdfParserAdapter(
                executable=sys.executable,
                method="auto",
                backend="pipeline",
                lang="ch",
                extra_args=[script_path],
            )
            blocks = adapter.parse_file(pdf_path)
        finally:
            os.remove(script_path)
            os.remove(pdf_path)

        self.assertEqual(blocks[0].text, "backend=pipeline")
        self.assertTrue(any(block.text == "lang=ch" for block in blocks))
        self.assertTrue(any(block.text == "method=auto" for block in blocks))

    def test_mineru_pdf_adapter_normalizes_trailing_whitespace_in_pdf_stem(self) -> None:
        adapter_script = """
from pathlib import Path
import sys

args = sys.argv[1:]
pdf_path = Path(args[args.index("-p") + 1])
output_dir = Path(args[args.index("-o") + 1])
if pdf_path.stem.endswith(" "):
    sys.stderr.write(f"bad-stem:{pdf_path.stem!r}\\n")
    raise SystemExit(3)
nested_dir = output_dir / pdf_path.stem / "auto"
nested_dir.mkdir(parents=True, exist_ok=True)
(nested_dir / f"{pdf_path.stem}.md").write_text("# 业务流程自动化\\n\\n解析成功", encoding="utf-8")
"""
        with tempfile.NamedTemporaryFile("w", suffix=".py", delete=False, encoding="utf-8") as script_handle:
            script_handle.write(adapter_script)
            script_path = script_handle.name

        with tempfile.TemporaryDirectory() as temp_dir:
            pdf_path = os.path.join(temp_dir, "业务流程自动化 Agent .pdf")
            with open(pdf_path, "wb") as pdf_handle:
                pdf_handle.write(self._build_valid_pdf_bytes())

            try:
                adapter = MinerUPdfParserAdapter(
                    executable=sys.executable,
                    extra_args=[script_path],
                )
                blocks = adapter.parse_file(pdf_path)
            finally:
                os.remove(script_path)

        self.assertEqual(blocks[0].text, "业务流程自动化")
        self.assertTrue(any(block.text == "解析成功" for block in blocks))

    def test_mineru_pdf_adapter_prefers_structured_json_and_keeps_page_numbers(self) -> None:
        adapter_script = """
from pathlib import Path
import json
import sys

args = sys.argv[1:]
pdf_path = Path(args[args.index("-p") + 1])
output_dir = Path(args[args.index("-o") + 1])
nested_dir = output_dir / pdf_path.stem / "auto"
nested_dir.mkdir(parents=True, exist_ok=True)
(nested_dir / f"{pdf_path.stem}.md").write_text("# should-not-win", encoding="utf-8")
(nested_dir / f"{pdf_path.stem}_content_list_v2.json").write_text(
    json.dumps(
        [
            [
                {
                    "type": "title",
                    "content": {
                        "title_content": [{"type": "text", "content": "价格指数管理办法"}],
                        "level": 1,
                    },
                    "bbox": [0, 0, 10, 10],
                },
                {
                    "type": "paragraph",
                    "content": {
                        "paragraph_content": [{"type": "text", "content": "第一章 总则"}]
                    },
                    "bbox": [0, 10, 10, 20],
                },
                {
                    "type": "paragraph",
                    "content": {
                        "paragraph_content": [{"type": "text", "content": "第一条 这里是正文内容，不是条级标题。"}]
                    },
                    "bbox": [0, 20, 10, 30],
                },
            ],
            [
                {
                    "type": "paragraph",
                    "content": {
                        "paragraph_content": [{"type": "text", "content": "第二条 这里是第二页正文。"}]
                    },
                    "bbox": [0, 0, 10, 10],
                }
            ],
        ],
        ensure_ascii=False,
    ),
    encoding="utf-8",
)
"""
        with tempfile.NamedTemporaryFile("w", suffix=".py", delete=False, encoding="utf-8") as script_handle:
            script_handle.write(adapter_script)
            script_path = script_handle.name

        with tempfile.NamedTemporaryFile("wb", suffix=".pdf", delete=False) as pdf_handle:
            pdf_handle.write(b"%PDF-1.4 fake")
            pdf_path = pdf_handle.name

        try:
            adapter = MinerUPdfParserAdapter(
                executable=sys.executable,
                extra_args=[script_path],
            )
            blocks = adapter.parse_file(pdf_path)
        finally:
            os.remove(script_path)
            os.remove(pdf_path)

        self.assertEqual(blocks[0].text, "价格指数管理办法")
        self.assertEqual(blocks[0].page_number, 1)
        self.assertEqual(blocks[1].block_type, "heading")
        self.assertEqual(blocks[1].text, "第一章 总则")
        self.assertEqual(blocks[1].page_number, 1)
        self.assertEqual(blocks[2].section_path, ["价格指数管理办法", "第一章 总则"])
        self.assertEqual(blocks[-1].page_number, 2)

    def test_mineru_pdf_adapter_can_parse_flat_content_list_json(self) -> None:
        adapter_script = """
from pathlib import Path
import json
import sys

args = sys.argv[1:]
pdf_path = Path(args[args.index("-p") + 1])
output_dir = Path(args[args.index("-o") + 1])
nested_dir = output_dir / pdf_path.stem / "auto"
nested_dir.mkdir(parents=True, exist_ok=True)
(nested_dir / f"{pdf_path.stem}_content_list.json").write_text(
    json.dumps(
        [
            {"type": "text", "text": "价格指数管理办法", "page_idx": 0, "text_level": 1},
            {"type": "text", "text": "第一章 总则", "page_idx": 0},
            {"type": "text", "text": "第一条 这是正文内容，不是标题。", "page_idx": 0},
            {"type": "text", "text": "第二章 行为主体", "page_idx": 2},
            {"type": "text", "text": "第六条 第二章正文内容。", "page_idx": 2},
        ],
        ensure_ascii=False,
    ),
    encoding="utf-8",
)
"""
        with tempfile.NamedTemporaryFile("w", suffix=".py", delete=False, encoding="utf-8") as script_handle:
            script_handle.write(adapter_script)
            script_path = script_handle.name

        with tempfile.NamedTemporaryFile("wb", suffix=".pdf", delete=False) as pdf_handle:
            pdf_handle.write(b"%PDF-1.4 fake")
            pdf_path = pdf_handle.name

        try:
            adapter = MinerUPdfParserAdapter(
                executable=sys.executable,
                extra_args=[script_path],
            )
            blocks = adapter.parse_file(pdf_path)
        finally:
            os.remove(script_path)
            os.remove(pdf_path)

        self.assertEqual(blocks[0].text, "价格指数管理办法")
        self.assertEqual(blocks[0].page_number, 1)
        self.assertEqual(blocks[1].block_type, "heading")
        self.assertEqual(blocks[1].text, "第一章 总则")
        self.assertEqual(blocks[2].section_path, ["价格指数管理办法", "第一章 总则"])
        self.assertEqual(blocks[-1].page_number, 3)

    def test_mineru_pdf_adapter_falls_back_to_markdown_when_structured_json_invalid(self) -> None:
        adapter_script = """
from pathlib import Path
import sys

args = sys.argv[1:]
pdf_path = Path(args[args.index("-p") + 1])
output_dir = Path(args[args.index("-o") + 1])
nested_dir = output_dir / pdf_path.stem / "auto"
nested_dir.mkdir(parents=True, exist_ok=True)
(nested_dir / f"{pdf_path.stem}_content_list_v2.json").write_text("{not-json", encoding="utf-8")
(nested_dir / f"{pdf_path.stem}.md").write_text("# 价格指数管理办法\\n\\n第一章 总则\\n\\n第一条 正文内容", encoding="utf-8")
"""
        with tempfile.NamedTemporaryFile("w", suffix=".py", delete=False, encoding="utf-8") as script_handle:
            script_handle.write(adapter_script)
            script_path = script_handle.name

        with tempfile.NamedTemporaryFile("wb", suffix=".pdf", delete=False) as pdf_handle:
            pdf_handle.write(b"%PDF-1.4 fake")
            pdf_path = pdf_handle.name

        try:
            adapter = MinerUPdfParserAdapter(
                executable=sys.executable,
                extra_args=[script_path],
            )
            blocks = adapter.parse_file(pdf_path)
        finally:
            os.remove(script_path)
            os.remove(pdf_path)

        self.assertEqual(blocks[0].block_type, "heading")
        self.assertEqual(blocks[0].text, "价格指数管理办法")
        self.assertTrue(any(block.text.startswith("第一条") for block in blocks))

    def test_mineru_pdf_adapter_keeps_article_titles_under_current_chapter(self) -> None:
        adapter_script = """
from pathlib import Path
import json
import sys

args = sys.argv[1:]
pdf_path = Path(args[args.index("-p") + 1])
output_dir = Path(args[args.index("-o") + 1])
nested_dir = output_dir / pdf_path.stem / "auto"
nested_dir.mkdir(parents=True, exist_ok=True)
(nested_dir / f"{pdf_path.stem}_content_list_v2.json").write_text(
    json.dumps(
        [
            [
                {
                    "type": "title",
                    "content": {
                        "title_content": [{"type": "text", "content": "价格指数管理办法"}],
                        "level": 1,
                    },
                },
                {
                    "type": "paragraph",
                    "content": {
                        "paragraph_content": [{"type": "text", "content": "第三章 价格指数的编制方案"}]
                    },
                },
                {
                    "type": "title",
                    "content": {
                        "title_content": [{"type": "text", "content": "第十一条 价格指数编制方案中保证价格指数完整性和可靠性措施包括："}],
                        "level": 2,
                    },
                },
                {
                    "type": "paragraph",
                    "content": {
                        "paragraph_content": [{"type": "text", "content": "（一）价格信息采集点的选择标准。"}]
                    },
                },
            ]
        ],
        ensure_ascii=False,
    ),
    encoding="utf-8",
)
"""
        with tempfile.NamedTemporaryFile("w", suffix=".py", delete=False, encoding="utf-8") as script_handle:
            script_handle.write(adapter_script)
            script_path = script_handle.name

        with tempfile.NamedTemporaryFile("wb", suffix=".pdf", delete=False) as pdf_handle:
            pdf_handle.write(b"%PDF-1.4 fake")
            pdf_path = pdf_handle.name

        try:
            adapter = MinerUPdfParserAdapter(
                executable=sys.executable,
                extra_args=[script_path],
            )
            blocks = adapter.parse_file(pdf_path)
        finally:
            os.remove(script_path)
            os.remove(pdf_path)

        self.assertEqual(blocks[2].block_type, "paragraph")
        self.assertEqual(
            blocks[2].section_path,
            ["价格指数管理办法", "第三章 价格指数的编制方案"],
        )
        self.assertEqual(
            blocks[3].section_path,
            ["价格指数管理办法", "第三章 价格指数的编制方案"],
        )

    def test_mineru_pdf_adapter_keeps_short_article_titles_as_subsections(self) -> None:
        adapter_script = """
from pathlib import Path
import json
import sys

args = sys.argv[1:]
pdf_path = Path(args[args.index("-p") + 1])
output_dir = Path(args[args.index("-o") + 1])
nested_dir = output_dir / pdf_path.stem / "auto"
nested_dir.mkdir(parents=True, exist_ok=True)
(nested_dir / f"{pdf_path.stem}_content_list_v2.json").write_text(
    json.dumps(
        [
            [
                {
                    "type": "title",
                    "content": {
                        "title_content": [{"type": "text", "content": "价格指数管理办法"}],
                        "level": 1,
                    },
                },
                {
                    "type": "paragraph",
                    "content": {
                        "paragraph_content": [{"type": "text", "content": "第三章 价格指数的编制方案"}]
                    },
                },
                {
                    "type": "title",
                    "content": {
                        "title_content": [{"type": "text", "content": "第十一条 定义"}],
                        "level": 2,
                    },
                },
                {
                    "type": "paragraph",
                    "content": {
                        "paragraph_content": [{"type": "text", "content": "这里是条级正文。"}]
                    },
                },
            ]
        ],
        ensure_ascii=False,
    ),
    encoding="utf-8",
)
"""
        with tempfile.NamedTemporaryFile("w", suffix=".py", delete=False, encoding="utf-8") as script_handle:
            script_handle.write(adapter_script)
            script_path = script_handle.name

        with tempfile.NamedTemporaryFile("wb", suffix=".pdf", delete=False) as pdf_handle:
            pdf_handle.write(b"%PDF-1.4 fake")
            pdf_path = pdf_handle.name

        try:
            adapter = MinerUPdfParserAdapter(
                executable=sys.executable,
                extra_args=[script_path],
            )
            blocks = adapter.parse_file(pdf_path)
        finally:
            os.remove(script_path)
            os.remove(pdf_path)

        self.assertEqual(blocks[2].block_type, "heading")
        self.assertEqual(
            blocks[2].section_path,
            ["价格指数管理办法", "第三章 价格指数的编制方案", "第十一条 定义"],
        )
        self.assertEqual(
            blocks[3].section_path,
            ["价格指数管理办法", "第三章 价格指数的编制方案", "第十一条 定义"],
        )

    def test_mineru_pdf_adapter_raises_clear_error_when_executable_missing(self) -> None:
        adapter = MinerUPdfParserAdapter(executable="missing-mineru-binary")

        with tempfile.NamedTemporaryFile("wb", suffix=".pdf", delete=False) as pdf_handle:
            pdf_handle.write(b"%PDF-1.4 fake")
            pdf_path = pdf_handle.name

        try:
            with self.assertRaises(ValueError) as context:
                adapter.parse_file(pdf_path)
        finally:
            os.remove(pdf_path)

        self.assertEqual(str(context.exception), "mineru_executable_not_found")

    def test_mineru_pdf_adapter_surfaces_cli_stderr_on_failure(self) -> None:
        adapter_script = """
import sys

sys.stderr.write("pipeline exploded\\n")
raise SystemExit(7)
"""
        with tempfile.NamedTemporaryFile("w", suffix=".py", delete=False, encoding="utf-8") as script_handle:
            script_handle.write(adapter_script)
            script_path = script_handle.name

        with tempfile.NamedTemporaryFile("wb", suffix=".pdf", delete=False) as pdf_handle:
            pdf_handle.write(b"%PDF-1.4 fake")
            pdf_path = pdf_handle.name

        try:
            adapter = MinerUPdfParserAdapter(
                executable=sys.executable,
                extra_args=[script_path],
            )
            with self.assertRaises(ValueError) as context:
                adapter.parse_file(pdf_path)
        finally:
            os.remove(script_path)
            os.remove(pdf_path)

        self.assertEqual(str(context.exception), "parse_failed:7:pipeline exploded")

    def test_mineru_pdf_adapter_raises_timeout_error(self) -> None:
        adapter_script = """
import time

time.sleep(0.3)
"""
        with tempfile.NamedTemporaryFile("w", suffix=".py", delete=False, encoding="utf-8") as script_handle:
            script_handle.write(adapter_script)
            script_path = script_handle.name

        with tempfile.NamedTemporaryFile("wb", suffix=".pdf", delete=False) as pdf_handle:
            pdf_handle.write(b"%PDF-1.4 fake")
            pdf_path = pdf_handle.name

        try:
            adapter = MinerUPdfParserAdapter(
                executable=sys.executable,
                extra_args=[script_path],
                timeout_seconds=0.05,
            )
            with self.assertRaises(ValueError) as context:
                adapter.parse_file(pdf_path)
        finally:
            os.remove(script_path)
            os.remove(pdf_path)

        self.assertEqual(str(context.exception), "parse_timeout")

    def test_chunker_splits_by_heading_before_size_limit(self) -> None:
        engine = ChunkPolicyEngine(soft_limit_chars=20, hard_limit_chars=40)
        blocks = [
            CanonicalBlock(block_type="heading", text="第一章"),
            CanonicalBlock(block_type="paragraph", text="报销制度第一段。"),
            CanonicalBlock(block_type="heading", text="第二章"),
            CanonicalBlock(block_type="paragraph", text="审批流程第二段。"),
        ]

        chunks = engine.build_chunks(blocks, source_file_id="file-1", title="制度")

        self.assertEqual(len(chunks), 2)
        self.assertEqual(chunks[0].title, "制度")

    def test_chunker_keeps_table_as_separate_chunk(self) -> None:
        engine = ChunkPolicyEngine(soft_limit_chars=200, hard_limit_chars=300)
        blocks = [
            CanonicalBlock(block_type="paragraph", text="报销制度说明"),
            CanonicalBlock(block_type="table", text="项目|金额\n交通|100"),
        ]

        chunks = engine.build_chunks(blocks, source_file_id="file-2", title="表格测试")

        self.assertEqual(len(chunks), 2)
        self.assertEqual(chunks[1].block_type, "table")

    def test_chunker_does_not_merge_across_slides(self) -> None:
        engine = ChunkPolicyEngine()
        blocks = [
            CanonicalBlock(block_type="slide", text="第一页内容", slide_number=1),
            CanonicalBlock(block_type="slide", text="第二页内容", slide_number=2),
        ]

        chunks = engine.build_chunks(blocks, source_file_id="file-3", title="销售PPT")

        self.assertEqual(chunks[0].slide_number, 1)
        self.assertEqual(chunks[1].slide_number, 2)

    def test_chunker_keeps_sheet_name_in_metadata(self) -> None:
        engine = ChunkPolicyEngine()
        blocks = [CanonicalBlock(block_type="sheet", text="报销记录", sheet_name="Sheet1")]

        chunks = engine.build_chunks(blocks, source_file_id="file-4", title="财务表")

        self.assertEqual(chunks[0].sheet_name, "Sheet1")

    def test_chunker_splits_overlong_chinese_paragraph_by_token_budget(self) -> None:
        engine = ChunkPolicyEngine(
            soft_limit_chars=900,
            hard_limit_chars=1400,
            soft_limit_tokens=100,
            hard_limit_tokens=160,
        )
        blocks = [CanonicalBlock(block_type="paragraph", text="报销制度说明。" * 200)]

        chunks = engine.build_chunks(blocks, source_file_id="file-5", title="长段落")

        self.assertGreater(len(chunks), 1)
        self.assertTrue(all(chunk.token_count_estimate <= 160 for chunk in chunks))

    def test_chunker_splits_wide_table_with_header_repeated(self) -> None:
        engine = ChunkPolicyEngine(
            soft_limit_chars=300,
            hard_limit_chars=500,
            soft_limit_tokens=80,
            hard_limit_tokens=120,
        )
        rows = "\n".join([f"行{i}|100|财务" for i in range(50)])
        blocks = [CanonicalBlock(block_type="table", text=f"项目|金额|部门\n{rows}")]

        chunks = engine.build_chunks(blocks, source_file_id="file-6", title="宽表")

        self.assertGreater(len(chunks), 1)
        self.assertTrue(all("项目|金额|部门" in chunk.content for chunk in chunks))
        self.assertTrue(all(chunk.token_count_estimate <= 120 for chunk in chunks))

    def test_chunker_splits_long_slide_notes_without_crossing_slides(self) -> None:
        engine = ChunkPolicyEngine(
            soft_limit_chars=300,
            hard_limit_chars=500,
            soft_limit_tokens=80,
            hard_limit_tokens=120,
        )
        blocks = [CanonicalBlock(block_type="slide", text="备注。" * 120, slide_number=3)]

        chunks = engine.build_chunks(blocks, source_file_id="file-7", title="备注页")

        self.assertGreater(len(chunks), 1)
        self.assertTrue(all(chunk.slide_number == 3 for chunk in chunks))

    def test_chunker_prefers_soft_budget_when_current_chunk_is_large_enough(self) -> None:
        engine = ChunkPolicyEngine(
            soft_limit_chars=14,
            hard_limit_chars=100,
            soft_limit_tokens=20,
            hard_limit_tokens=100,
            min_merge_chars=6,
        )
        blocks = [
            CanonicalBlock(
                block_type="paragraph",
                text="第一句内容很长。第二句内容很长。第三句内容很长。",
            )
        ]

        chunks = engine.build_chunks(blocks, source_file_id="file-8", title="软阈值")

        self.assertGreater(len(chunks), 1)

    def test_chunker_inherits_section_path_from_heading_block(self) -> None:
        engine = ChunkPolicyEngine()
        blocks = [
            CanonicalBlock(block_type="heading", text="第一章", section_path=["第一章"]),
            CanonicalBlock(block_type="paragraph", text="报销制度正文"),
        ]

        chunks = engine.build_chunks(blocks, source_file_id="file-9", title="制度")

        self.assertEqual(chunks[0].section_path, ["第一章"])

    def test_chunker_uses_group_start_page_as_anchor(self) -> None:
        engine = ChunkPolicyEngine()
        blocks = [
            CanonicalBlock(block_type="heading", text="第一章", section_path=["第一章"], page_number=1),
            CanonicalBlock(block_type="paragraph", text="第一页正文。", section_path=["第一章"], page_number=1),
            CanonicalBlock(block_type="paragraph", text="第二页正文。", section_path=["第一章"], page_number=2),
        ]

        chunks = engine.build_chunks(blocks, source_file_id="file-10", title="制度")

        self.assertEqual(len(chunks), 1)
        self.assertEqual(chunks[0].page_number, 1)

    def test_chunker_splits_single_oversized_table_row_to_budget(self) -> None:
        engine = ChunkPolicyEngine(
            soft_limit_chars=120,
            hard_limit_chars=180,
            soft_limit_tokens=50,
            hard_limit_tokens=70,
        )
        long_cell = "超长备注" * 80
        blocks = [CanonicalBlock(block_type="table", text=f"项目|备注\n差旅|{long_cell}")]

        chunks = engine.build_chunks(blocks, source_file_id="file-10", title="超长表格")

        self.assertGreater(len(chunks), 1)
        self.assertTrue(all(chunk.token_count_estimate <= 70 for chunk in chunks))

    def test_chunker_keeps_hard_budget_when_table_header_itself_is_oversized(self) -> None:
        engine = ChunkPolicyEngine(
            soft_limit_chars=120,
            hard_limit_chars=180,
            soft_limit_tokens=50,
            hard_limit_tokens=70,
        )
        header = "列名" * 120
        blocks = [CanonicalBlock(block_type="table", text=f"{header}\n差旅|200")]

        chunks = engine.build_chunks(blocks, source_file_id="file-11", title="超长表头")

        self.assertTrue(chunks)
        self.assertTrue(all(len(chunk.content) <= 180 for chunk in chunks))
        self.assertTrue(all(chunk.token_count_estimate <= 70 for chunk in chunks))


if __name__ == "__main__":
    unittest.main()
