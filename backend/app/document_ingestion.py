import math
import json
import os
from pathlib import Path
import re
import shutil
import subprocess
import tempfile
from typing import Dict, Iterable, List, Optional, Tuple

from pypdf import PdfReader

from app.document_models import CanonicalBlock, KnowledgeChunkRecord


def _normalize_text(value) -> str:
    return str(value).strip() if value is not None else ""


class TxtParserAdapter:
    def parse_file(self, file_path: str) -> List[CanonicalBlock]:
        with open(file_path, "r", encoding="utf-8-sig") as handle:
            text = handle.read().strip()
        return [CanonicalBlock(block_type="paragraph", text=text)] if text else []


class DocxParserAdapter:
    def parse_file(self, file_path: str) -> List[CanonicalBlock]:
        from docx import Document

        document = Document(file_path)
        blocks: List[CanonicalBlock] = []
        section_path: List[str] = []

        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue
            style_name = paragraph.style.name if paragraph.style is not None else ""
            heading_level = self._heading_level(style_name)
            if heading_level is not None:
                section_path = section_path[: heading_level - 1] + [text]
                blocks.append(
                    CanonicalBlock(
                        block_type="heading",
                        text=text,
                        section_path=list(section_path),
                        metadata={"source_type": "docx", "heading_level": heading_level},
                    )
                )
                continue
            blocks.append(
                CanonicalBlock(
                    block_type="paragraph",
                    text=text,
                    section_path=list(section_path),
                    metadata={"source_type": "docx"},
                )
            )
        return blocks

    @staticmethod
    def _heading_level(style_name: str) -> Optional[int]:
        match = re.search(r"Heading\s+(\d+)", style_name or "", re.IGNORECASE)
        if not match:
            return None
        return int(match.group(1))


class PptxParserAdapter:
    def parse_file(self, file_path: str) -> List[CanonicalBlock]:
        from pptx import Presentation

        presentation = Presentation(file_path)
        blocks: List[CanonicalBlock] = []
        for index, slide in enumerate(presentation.slides, start=1):
            title_shape = slide.shapes.title
            slide_title = title_shape.text.strip() if title_shape and title_shape.text else ""
            section_path = [slide_title] if slide_title else []
            if slide_title:
                blocks.append(
                    CanonicalBlock(
                        block_type="heading",
                        text=slide_title,
                        section_path=list(section_path),
                        slide_number=index,
                        metadata={"source_type": "pptx"},
                    )
                )
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                text = shape.text.strip()
                if not text:
                    continue
                if title_shape is not None and shape == title_shape:
                    continue
                blocks.append(
                    CanonicalBlock(
                        block_type="slide",
                        text=text,
                        section_path=list(section_path),
                        slide_number=index,
                        metadata={"source_type": "pptx"},
                    )
                )
        return blocks


class XlsxParserAdapter:
    def parse_file(self, file_path: str) -> List[CanonicalBlock]:
        from openpyxl import load_workbook

        workbook = load_workbook(file_path, read_only=True, data_only=True)
        blocks: List[CanonicalBlock] = []
        for worksheet in workbook.worksheets:
            rows = [
                [_normalize_text(value) for value in row]
                for row in worksheet.iter_rows(values_only=True)
            ]
            rows = [row for row in rows if any(cell for cell in row)]
            if not rows:
                continue
            lines = [" | ".join(cell for cell in row if cell) for row in rows]
            content = "\n".join(line for line in lines if line)
            if not content:
                continue
            blocks.append(
                CanonicalBlock(
                    block_type="table",
                    text=content,
                    sheet_name=worksheet.title,
                    metadata={"source_type": "xlsx", "row_count": len(rows)},
                )
            )
        workbook.close()
        return blocks


class MinerUPdfParserAdapter:
    def __init__(
        self,
        executable: Optional[str] = None,
        method: Optional[str] = None,
        backend: Optional[str] = None,
        lang: Optional[str] = None,
        timeout_seconds: Optional[float] = None,
        extra_args: Optional[List[str]] = None,
    ) -> None:
        self.executable = executable or os.getenv("EMATA_MINERU_EXECUTABLE", "mineru")
        self.method = method or os.getenv("EMATA_MINERU_METHOD", "auto")
        self.backend = backend or os.getenv("EMATA_MINERU_BACKEND", "pipeline")
        self.lang = lang or os.getenv("EMATA_MINERU_LANG", "ch")
        self.timeout_seconds = float(timeout_seconds or os.getenv("EMATA_MINERU_TIMEOUT_SECONDS", "600"))
        self.extra_args = list(extra_args or [])

    def parse_file(self, file_path: str) -> List[CanonicalBlock]:
        try:
            self._validate_pdf_file(
                file_path=file_path,
                strict_validation=Path(self.executable).name.lower().startswith("mineru"),
            )
        except FileNotFoundError as exc:
            if not self._is_executable_available():
                raise ValueError("mineru_executable_not_found") from exc
            raise ValueError("invalid_pdf_file") from exc
        if not self._is_executable_available():
            raise ValueError("mineru_executable_not_found")
        with tempfile.TemporaryDirectory(prefix="emata-mineru-") as output_dir:
            prepared_file_path = self._prepare_mineru_input_path(
                file_path=file_path,
                working_dir=Path(output_dir),
            )
            pdf_stem = prepared_file_path.stem
            command = [
                self.executable,
                *self.extra_args,
                "-p",
                str(prepared_file_path),
                "-o",
                output_dir,
                "-m",
                self.method,
            ]
            if self.backend:
                command.extend(["-b", self.backend])
            if self.lang:
                command.extend(["-l", self.lang])
            try:
                subprocess.run(
                    command,
                    check=True,
                    capture_output=True,
                    text=True,
                    timeout=self.timeout_seconds,
                )
            except FileNotFoundError as exc:
                raise ValueError("mineru_executable_not_found") from exc
            except subprocess.TimeoutExpired as exc:
                raise ValueError("parse_timeout") from exc
            except subprocess.CalledProcessError as exc:
                summary = self._summarize_process_error(exc)
                detail = f"parse_failed:{exc.returncode}"
                if summary:
                    detail = f"{detail}:{summary}"
                raise ValueError(detail) from exc

            structured_json_path = self._resolve_structured_json_path(
                output_dir=Path(output_dir),
                pdf_stem=pdf_stem,
            )
            if structured_json_path is not None:
                try:
                    blocks = self._parse_structured_json(
                        json.loads(structured_json_path.read_text(encoding="utf-8"))
                    )
                    if blocks:
                        return blocks
                except (json.JSONDecodeError, TypeError, ValueError, AttributeError):
                    pass

            markdown_path = self._resolve_markdown_path(
                output_dir=Path(output_dir),
                pdf_stem=pdf_stem,
            )
            if markdown_path is None:
                raise ValueError("mineru_output_missing")
            return self._parse_markdown(markdown_path.read_text(encoding="utf-8"))

    def _is_executable_available(self) -> bool:
        executable = (self.executable or "").strip()
        if not executable:
            return False
        if os.path.isabs(executable) or os.path.sep in executable:
            return Path(executable).exists()
        return shutil.which(executable) is not None

    @staticmethod
    def _validate_pdf_file(file_path: str, *, strict_validation: bool = False) -> None:
        try:
            with open(file_path, "rb") as handle:
                header = handle.read(5)
            if header != b"%PDF-":
                raise ValueError("invalid_pdf_file")
            if strict_validation:
                reader = PdfReader(file_path)
                if len(reader.pages) == 0:
                    raise ValueError("invalid_pdf_file")
        except ValueError:
            raise
        except FileNotFoundError:
            raise
        except Exception as exc:
            raise ValueError("invalid_pdf_file") from exc

    @staticmethod
    def _prepare_mineru_input_path(file_path: str, working_dir: Path) -> Path:
        source_path = Path(file_path)
        normalized_stem = source_path.stem.strip()
        if normalized_stem == source_path.stem:
            return source_path

        normalized_name = f"{normalized_stem}{source_path.suffix}"
        normalized_path = working_dir / normalized_name
        shutil.copy2(source_path, normalized_path)
        return normalized_path

    def _resolve_markdown_path(self, output_dir: Path, pdf_stem: str) -> Optional[Path]:
        candidates = [
            output_dir / f"{pdf_stem}.md",
            output_dir / pdf_stem / self.method / f"{pdf_stem}.md",
        ]
        for candidate in candidates:
            if candidate.exists():
                return candidate

        for candidate in output_dir.rglob(f"{pdf_stem}.md"):
            if candidate.is_file():
                return candidate
        return None

    def _resolve_structured_json_path(self, output_dir: Path, pdf_stem: str) -> Optional[Path]:
        candidates = [
            output_dir / f"{pdf_stem}_content_list_v2.json",
            output_dir / pdf_stem / self.method / f"{pdf_stem}_content_list_v2.json",
            output_dir / f"{pdf_stem}_content_list.json",
            output_dir / pdf_stem / self.method / f"{pdf_stem}_content_list.json",
        ]
        for candidate in candidates:
            if candidate.exists():
                return candidate
        for suffix in ("_content_list_v2.json", "_content_list.json"):
            for candidate in output_dir.rglob(f"{pdf_stem}{suffix}"):
                if candidate.is_file():
                    return candidate
        return None

    def _parse_structured_json(self, payload) -> List[CanonicalBlock]:
        blocks: List[CanonicalBlock] = []
        section_path: List[str] = []

        for page_number, page_blocks in self._iterate_structured_pages(payload):
            for block in page_blocks:
                block_type = block.get("type")
                text = self._extract_structured_text(block).strip()
                if not text or block_type == "page_number":
                    continue

                raw_title_level = None
                if block_type == "title":
                    raw_title_level = block.get("content", {}).get("level", 1)
                elif block_type == "text" and block.get("text_level") is not None:
                    raw_title_level = block.get("text_level")

                if raw_title_level is not None:
                    level = int(raw_title_level or 1)
                    if self._should_demote_article_title(text):
                        blocks.append(
                            CanonicalBlock(
                                block_type="paragraph",
                                text=text,
                                section_path=list(section_path),
                                page_number=page_number,
                                metadata={"source_type": "pdf", "parser": "mineru"},
                            )
                        )
                        continue
                    if self._should_keep_article_title_as_heading(text):
                        level = max(level, 3)
                    section_path = section_path[: level - 1] + [text]
                    blocks.append(
                        CanonicalBlock(
                            block_type="heading",
                            text=text,
                            section_path=list(section_path),
                            page_number=page_number,
                            metadata={
                                "source_type": "pdf",
                                "parser": "mineru",
                                "heading_level": level,
                            },
                        )
                    )
                    continue

                heading_level = self._infer_structured_heading_level(text)
                if heading_level is not None:
                    section_path = section_path[: heading_level - 1] + [text]
                    blocks.append(
                        CanonicalBlock(
                            block_type="heading",
                            text=text,
                            section_path=list(section_path),
                            page_number=page_number,
                            metadata={
                                "source_type": "pdf",
                                "parser": "mineru",
                                "heading_level": heading_level,
                            },
                        )
                    )
                    continue

                normalized_block_type = "table" if block_type == "table" else "paragraph"
                blocks.append(
                    CanonicalBlock(
                        block_type=normalized_block_type,
                        text=text,
                        section_path=list(section_path),
                        page_number=page_number,
                        metadata={"source_type": "pdf", "parser": "mineru"},
                    )
                )
        return blocks

    def _iterate_structured_pages(self, payload) -> Iterable[Tuple[int, List[dict]]]:
        if not isinstance(payload, list):
            raise ValueError("structured_json_invalid")
        if not payload:
            return []

        if all(isinstance(item, list) for item in payload):
            return [
                (page_number, page_blocks)
                for page_number, page_blocks in enumerate(payload, start=1)
            ]

        if all(isinstance(item, dict) for item in payload):
            grouped: Dict[int, List[dict]] = {}
            for block in payload:
                page_number = self._resolve_structured_page_number(block)
                grouped.setdefault(page_number, []).append(block)
            return [(page_number, grouped[page_number]) for page_number in sorted(grouped)]

        raise ValueError("structured_json_invalid")

    @staticmethod
    def _resolve_structured_page_number(block: dict) -> int:
        if "page_idx" in block and block.get("page_idx") is not None:
            return int(block["page_idx"]) + 1
        for key in ("page_number", "page_no", "page"):
            if block.get(key) is not None:
                return int(block[key])
        return 1

    def _should_demote_article_title(self, text: str) -> bool:
        if not re.match(r"^第[一二三四五六七八九十百千万零〇0-9]+条", text):
            return False
        if self._should_keep_article_title_as_heading(text):
            return False
        return True

    @staticmethod
    def _should_keep_article_title_as_heading(text: str) -> bool:
        match = re.match(r"^第[一二三四五六七八九十百千万零〇0-9]+条\s*(.*)$", text)
        if not match:
            return False
        tail = match.group(1).strip()
        if not tail:
            return False
        if any(marker in tail for marker in ("：", ":", "；", ";", "。", "，", ",")):
            return False
        return len(tail) <= 12

    def _parse_markdown(self, content: str) -> List[CanonicalBlock]:
        blocks: List[CanonicalBlock] = []
        section_path: List[str] = []
        paragraph_lines: List[str] = []
        table_lines: List[str] = []

        def flush_paragraph() -> None:
            nonlocal paragraph_lines
            text = "\n".join(paragraph_lines).strip()
            if text:
                blocks.append(
                    CanonicalBlock(
                        block_type="paragraph",
                        text=text,
                        section_path=list(section_path),
                        metadata={"source_type": "pdf", "parser": "mineru"},
                    )
                )
            paragraph_lines = []

        def flush_table() -> None:
            nonlocal table_lines
            text = "\n".join(table_lines).strip()
            if text:
                blocks.append(
                    CanonicalBlock(
                        block_type="table",
                        text=text,
                        section_path=list(section_path),
                        metadata={"source_type": "pdf", "parser": "mineru"},
                    )
                )
            table_lines = []

        for raw_line in content.splitlines():
            line = raw_line.strip()
            if not line:
                flush_paragraph()
                flush_table()
                continue
            if line.startswith("#"):
                flush_paragraph()
                flush_table()
                level = len(line) - len(line.lstrip("#"))
                title = line[level:].strip()
                section_path = section_path[: level - 1] + [title]
                blocks.append(
                    CanonicalBlock(
                        block_type="heading",
                        text=title,
                        section_path=list(section_path),
                        metadata={
                            "source_type": "pdf",
                            "parser": "mineru",
                            "heading_level": level,
                        },
                    )
                )
                continue
            if "|" in line:
                flush_paragraph()
                table_lines.append(line)
                continue
            flush_table()
            paragraph_lines.append(line)

        flush_paragraph()
        flush_table()
        return blocks

    @staticmethod
    def _extract_structured_text(block: dict) -> str:
        raw_text = block.get("text")
        if isinstance(raw_text, str) and raw_text.strip():
            return raw_text.strip()
        content = block.get("content", {})
        keys = ("title_content", "paragraph_content", "table_caption", "table_body", "text")
        parts: List[str] = []
        for key in keys:
            value = content.get(key)
            if not value:
                continue
            if isinstance(value, str):
                parts.append(value)
                continue
            if isinstance(value, list):
                for item in value:
                    if isinstance(item, dict):
                        nested_content = item.get("content")
                        if isinstance(nested_content, str):
                            parts.append(nested_content)
                    elif isinstance(item, str):
                        parts.append(item)
        return "\n".join(part.strip() for part in parts if part and part.strip())

    @staticmethod
    def _infer_structured_heading_level(text: str) -> Optional[int]:
        if re.match(r"^第[一二三四五六七八九十百千万零〇0-9]+章\b", text):
            return 2
        if MinerUPdfParserAdapter._should_keep_article_title_as_heading(text):
            return 3
        return None

    @staticmethod
    def _summarize_process_error(exc: subprocess.CalledProcessError) -> str:
        combined = "\n".join(part for part in [exc.stderr, exc.stdout] if part).strip()
        if not combined:
            return ""

        for line in reversed([item.strip() for item in combined.splitlines() if item.strip()]):
            error_match = re.search(r'"error":\s*"([^"]+)"', line)
            if error_match:
                return MinerUPdfParserAdapter._truncate_error(error_match.group(1))
            if "Error:" in line:
                return MinerUPdfParserAdapter._truncate_error(line)
            if "No module named" in line:
                return MinerUPdfParserAdapter._truncate_error(line)
            if line.startswith("INFO:") or line.startswith("WARNING:"):
                continue
            return MinerUPdfParserAdapter._truncate_error(line)
        return ""

    @staticmethod
    def _truncate_error(message: str) -> str:
        normalized = " ".join(message.split())
        return normalized[:240]


class DocumentParserRegistry:
    def __init__(self) -> None:
        self._parsers: Dict[str, object] = {
            "txt": TxtParserAdapter(),
            "docx": DocxParserAdapter(),
            "pptx": PptxParserAdapter(),
            "xlsx": XlsxParserAdapter(),
            "pdf": MinerUPdfParserAdapter(),
        }

    def parse_file(self, file_path: str, source_type: str) -> List[CanonicalBlock]:
        parser = self._parsers.get(source_type)
        if parser is None:
            raise ValueError("unsupported_source_type")
        return parser.parse_file(file_path)

    def _normalize_mock_blocks(self, raw_blocks: Iterable[dict]) -> List[CanonicalBlock]:
        return [
            CanonicalBlock(
                block_type=item["type"],
                text=item["text"],
                section_path=item.get("section_path", []),
                page_number=item.get("page_number"),
                sheet_name=item.get("sheet_name"),
                slide_number=item.get("slide_number"),
                metadata=item.get("metadata", {}),
            )
            for item in raw_blocks
        ]


class DoclingParserAdapter:
    def __init__(self) -> None:
        self.registry = DocumentParserRegistry()

    def parse_file(self, file_path: str, source_type: str) -> List[CanonicalBlock]:
        return self.registry.parse_file(file_path, source_type)

    def _normalize_mock_blocks(self, raw_blocks: Iterable[dict]) -> List[CanonicalBlock]:
        return self.registry._normalize_mock_blocks(raw_blocks)


class ChunkPolicyEngine:
    def __init__(
        self,
        soft_limit_chars: int = 900,
        hard_limit_chars: int = 1400,
        soft_limit_tokens: int = 600,
        hard_limit_tokens: int = 900,
        min_merge_chars: int = 220,
    ) -> None:
        self.soft_limit_chars = soft_limit_chars
        self.hard_limit_chars = hard_limit_chars
        self.soft_limit_tokens = soft_limit_tokens
        self.hard_limit_tokens = hard_limit_tokens
        self.min_merge_chars = min_merge_chars

    def build_chunks(
        self,
        blocks: List[CanonicalBlock],
        source_file_id: str,
        title: str,
        organization_id: str = "",
        workspace_id: Optional[str] = None,
        scope: str = "workspace",
    ) -> List[KnowledgeChunkRecord]:
        if not blocks:
            return []

        groups: List[List[CanonicalBlock]] = []
        current: List[CanonicalBlock] = []
        last_slide: Optional[int] = None
        last_sheet: Optional[str] = None

        for block in blocks:
            slide_boundary = block.slide_number is not None and block.slide_number != last_slide
            sheet_boundary = block.sheet_name is not None and block.sheet_name != last_sheet
            heading_boundary = block.block_type == "heading" and len(current) > 0
            table_boundary = block.block_type == "table"

            if table_boundary:
                if current:
                    groups.append(current)
                    current = []
                groups.append([block])
            else:
                if (slide_boundary or sheet_boundary or heading_boundary) and current:
                    groups.append(current)
                    current = []
                current.append(block)

            if block.slide_number is not None:
                last_slide = block.slide_number
            if block.sheet_name is not None:
                last_sheet = block.sheet_name

        if current:
            groups.append(current)

        chunks: List[KnowledgeChunkRecord] = []
        for group in groups:
            chunks.extend(
                self._build_group_chunks(
                    group=group,
                    source_file_id=source_file_id,
                    title=title,
                    organization_id=organization_id,
                    workspace_id=workspace_id,
                    scope=scope,
                    start_index=len(chunks),
                )
            )
        return chunks

    def _build_group_chunks(
        self,
        group: List[CanonicalBlock],
        source_file_id: str,
        title: str,
        organization_id: str,
        workspace_id: Optional[str],
        scope: str,
        start_index: int,
    ) -> List[KnowledgeChunkRecord]:
        block_type = self._resolve_block_type(group)
        anchor = self._resolve_anchor(group)
        page_numbers = [block.page_number for block in group if block.page_number is not None]
        page_end = max(page_numbers) if page_numbers else None
        metadata = self._resolve_chunk_metadata(group, page_end)
        if block_type == "table":
            texts = self._split_table_text(group[0].text)
        else:
            merged_text = "\n".join(block.text.strip() for block in group if block.text.strip())
            texts = self._split_text_by_budget(merged_text)

        built: List[KnowledgeChunkRecord] = []
        for offset, text in enumerate(texts):
            built.append(
                KnowledgeChunkRecord(
                    id=f"{source_file_id}-chunk-{start_index + offset}",
                    source_file_id=source_file_id,
                    organization_id=organization_id,
                    workspace_id=workspace_id,
                    scope=scope,
                    title=title,
                    content=text,
                    block_type=block_type,
                    section_path=anchor.section_path,
                    page_number=anchor.page_number,
                    sheet_name=anchor.sheet_name,
                    slide_number=anchor.slide_number,
                    chunk_index=start_index + offset,
                    token_count_estimate=self._estimate_tokens(text),
                    metadata=dict(metadata),
                )
            )
        return built

    def _resolve_block_type(self, group: List[CanonicalBlock]) -> str:
        for block in reversed(group):
            if block.block_type != "heading":
                return block.block_type
        return group[-1].block_type

    def _resolve_anchor(self, group: List[CanonicalBlock]) -> CanonicalBlock:
        section_path = []
        page_number = None
        sheet_name = None
        slide_number = None
        for block in group:
            if not section_path and block.section_path:
                section_path = block.section_path
            if page_number is None and block.page_number is not None:
                page_number = block.page_number
            if sheet_name is None and block.sheet_name is not None:
                sheet_name = block.sheet_name
            if slide_number is None and block.slide_number is not None:
                slide_number = block.slide_number
        return CanonicalBlock(
            block_type=self._resolve_block_type(group),
            text="",
            section_path=section_path,
            page_number=page_number,
            sheet_name=sheet_name,
            slide_number=slide_number,
        )

    @staticmethod
    def _resolve_chunk_metadata(group: List[CanonicalBlock], page_end: Optional[int]) -> Dict[str, object]:
        metadata: Dict[str, object] = {}
        for block in group:
            if not block.metadata:
                continue
            for key, value in block.metadata.items():
                metadata.setdefault(key, value)
        if page_end is not None:
            metadata["page_end"] = page_end
        return metadata

    def _split_text_by_budget(self, text: str) -> List[str]:
        text = text.strip()
        if not text:
            return [""]
        if self._fits_soft_budget(text):
            return [text]

        chunks: List[str] = []
        current = ""
        for sentence in self._split_sentences(text):
            candidate = f"{current}{sentence}" if current else sentence
            if self._fits_soft_budget(candidate):
                current = candidate
                continue

            if current and self._fits_budget(candidate) and len(current) < self.min_merge_chars:
                current = candidate
                continue

            if current and self._fits_budget(candidate):
                chunks.append(current.strip())
                current = sentence
                continue

            if current:
                chunks.append(current.strip())

            if self._fits_budget(sentence):
                current = sentence
            else:
                chunks.extend(self._hard_split(sentence))
                current = ""

        if current:
            chunks.append(current.strip())
        return [chunk for chunk in chunks if chunk]

    def _split_table_text(self, text: str) -> List[str]:
        lines = [line for line in text.splitlines() if line.strip()]
        if len(lines) <= 1:
            return self._hard_split(text.strip())

        header = lines[0]
        if not self._fits_budget(header):
            return self._hard_split(text.strip())
        rows = lines[1:]
        chunks: List[str] = []
        current_rows: List[str] = []

        for row in rows:
            candidate_rows = current_rows + [row]
            candidate_text = "\n".join([header] + candidate_rows)
            if self._fits_soft_budget(candidate_text):
                current_rows = candidate_rows
                continue

            if current_rows and self._fits_budget(candidate_text) and len("\n".join(current_rows)) < self.min_merge_chars:
                current_rows = candidate_rows
                continue

            if current_rows:
                chunks.append("\n".join([header] + current_rows))
                current_rows = []

            if self._fits_budget("\n".join([header, row])):
                current_rows = [row]
                continue

            allowed_chars = max(1, min(self.hard_limit_chars, self.hard_limit_tokens * 2) - len(header) - 1)
            split_row_chunks = self._hard_split(row, max_chars_override=allowed_chars)
            for piece in split_row_chunks:
                chunks.append("\n".join([header, piece]))

        if current_rows:
            chunks.append("\n".join([header] + current_rows))
        return chunks

    def _split_sentences(self, text: str) -> List[str]:
        parts = re.split(r"(?<=[。！？；.!?;\n])", text)
        return [part for part in parts if part and part.strip()]

    def _hard_split(self, text: str, max_chars_override: Optional[int] = None) -> List[str]:
        pieces: List[str] = []
        start = 0
        max_chars = max_chars_override or min(self.hard_limit_chars, self.hard_limit_tokens * 2)
        max_chars = max(1, max_chars)
        while start < len(text):
            end = min(start + max_chars, len(text))
            pieces.append(text[start:end].strip())
            start = end
        return [piece for piece in pieces if piece]

    def _estimate_tokens(self, text: str) -> int:
        return max(1, math.ceil(len(text) / 2))

    def _fits_budget(self, text: str) -> bool:
        return len(text) <= self.hard_limit_chars and self._estimate_tokens(text) <= self.hard_limit_tokens

    def _fits_soft_budget(self, text: str) -> bool:
        return len(text) <= self.soft_limit_chars and self._estimate_tokens(text) <= self.soft_limit_tokens

    def _needs_secondary_split(self, text: str) -> bool:
        return len(text) > self.hard_limit_chars or self._estimate_tokens(text) > self.hard_limit_tokens
